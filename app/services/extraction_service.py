"""
Text extraction from uploaded files.

Supports: PDF, XLSX, XLSM, DOCX, PPTX, TXT.
Takes raw bytes + filename; returns extracted plain text.
"""
from __future__ import annotations

import io
import zipfile

ALLOWED_EXTENSIONS = {"pdf", "xlsx", "xlsm", "docx", "pptx", "txt"}
ALLOWED_CONTENT_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel.sheet.macroEnabled.12",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "text/plain",
    # Browsers sometimes send generic types for Office files
    "application/octet-stream",
    "application/zip",
}


def get_extension(filename: str) -> str:
    """Return lowercase extension without the dot, e.g. 'pdf'."""
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def is_allowed(filename: str) -> bool:
    return get_extension(filename) in ALLOWED_EXTENSIONS


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------

def _extract_pdf(data: bytes) -> str:
    import pdfplumber
    text_parts: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
    return "\n\n".join(text_parts)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    # Paragraphs
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_txt(data: bytes) -> str:
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


# ---------------------------------------------------------------------------
# Public interface
# ---------------------------------------------------------------------------

def extract_text(filename: str, data: bytes) -> str:
    """
    Extract plain text from file bytes.
    Raises ValueError for unsupported file types.
    Raises RuntimeError if extraction fails.
    """
    ext = get_extension(filename)
    try:
        if ext == "pdf":
            return _extract_pdf(data)
        elif ext in ("xlsx", "xlsm"):
            return _extract_xlsx(data)
        elif ext == "docx":
            return _extract_docx(data)
        elif ext == "pptx":
            return _extract_pptx(data)
        elif ext == "txt":
            return _extract_txt(data)
        else:
            raise ValueError(f"Unsupported file type: .{ext}")
    except ValueError:
        raise
    except Exception as exc:
        raise RuntimeError(f"Text extraction failed for {filename!r}: {exc}") from exc
